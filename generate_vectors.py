import os
import logging
from llama_index.llms.openai import OpenAI
from llama_index.core import (
    SimpleDirectoryReader,
    VectorStoreIndex,
    StorageContext,
    load_index_from_storage,
)
from llama_index.core.schema import Document

from llama_index.core.tools import QueryEngineTool, ToolMetadata
from llama_index.core import Settings


from pptx import Presentation
import fitz

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")
logger = logging.getLogger(__name__)

# os.environ['OPENAI_API_KEY'] = os.getenv('API_KEY')
os.environ['OPENAI_API_KEY'] = ""

VECTOR_STORAGE_DIR = "./vector"
DATA_STORAGE_DIR = "./data"

os.makedirs(DATA_STORAGE_DIR, exist_ok=True)

def load_pptx_text(file_path):
    """Given a pptx file path, extract all the text from it"""
    pres = Presentation(file_path)
    pptx_text = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        pptx_text.append(run.text)
                        
    return pptx_text

def extract_text_from_pdf(file_path):
    """Extract text from a PDF file using PyMuPDF."""
    text = []
    try:
        pdf_document = fitz.open(file_path)
        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]
            text.append(page.get_text())
        pdf_document.close()
    except Exception as e:
        print(f"Error reading PDF file {file_path}: {e}")
    return "\n".join(text)

def load_any_documents(dirpath:str):
    """Given a directory path, load all the files in it: .pptx, .docx, .txt, .pdf, png"""
    documents = []
    for file_name in os.listdir(dirpath):
        file_path = os.path.join(dirpath, file_name)
        
        if file_name.endswith(".pptx"):
            # handle pptx file
            pptx_text = load_pptx_text(file_path)
            documents.append(Document(text="\n".join(pptx_text), doc_id=file_name))
        elif file_name.endswith(".pdf"):
            pdf_text = extract_text_from_pdf(file_path)
            documents.append(Document(text=pdf_text, doc_id=file_name))
        else:
            documents += (SimpleDirectoryReader(input_files=[file_path]).load_data())
            
    return documents

# initialize an index list
index_list = []
# description_list = []
## represent index represent a subdirectory

os.makedirs(VECTOR_STORAGE_DIR, exist_ok=True)

try:
    storage_context = StorageContext.from_defaults(persist_dir=VECTOR_STORAGE_DIR)
    index = load_index_from_storage(storage_context)

    index_list.append(index)

except FileNotFoundError:
    logger.info(f"Index for main directory not found. Creating a new index.")

    documents = load_any_documents(DATA_STORAGE_DIR)
    index = VectorStoreIndex.from_documents(documents)
    index.storage_context.persist(VECTOR_STORAGE_DIR)
    index_list.append(index)

except Exception as e:
    logger.error(f"Error processing main directory: {e}")   